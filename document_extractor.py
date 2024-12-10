import ray
from pathlib import Path
import pdfplumber
from pptx import Presentation
from pptx.exc import PackageNotFoundError

def extract_pdf(pdf_path):
    pages = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                pages.append({
                    "page_number": page_number,
                    "text": text
                })
    except Exception as e:
        print(f"Errore nell'estrazione del testo da {pdf_path}: {e}")
    return pages

def extract_ppt(ppt_path):
    slides = []
    try:
        prs = Presentation(ppt_path)
        for slide_number, slide in enumerate(prs.slides, start=1):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            slides.append({
                "page_number": slide_number,
                "text": text
            })
    except (PackageNotFoundError, KeyError) as e:
        print(f"Errore nell'apertura del file {ppt_path}: {e}")
    except Exception as e:
        print(f"Errore generico nell'elaborazione di {ppt_path}: {e}")
    return slides

# RAY per processare i documenti in parallelo (qualora servisse)
@ray.remote
def process_document(path):
    content = []
    path = Path(path)
    if path.suffix == '.pdf':
        content = extract_pdf(path)
    elif path.suffix in ['.ppt', '.pptx']:
        content = extract_ppt(path)
    return {"path": str(path), "content": content}

class DocumentExtractor:
    def __init__(self, docs_dir):
        self.docs_dir = Path(docs_dir)

    def get_documents(self):
        if self.docs_dir.is_file():
            paths = [self.docs_dir]
        elif self.docs_dir.is_dir():
            paths = [path for path in self.docs_dir.rglob("*") if path.suffix in ['.pdf', '.ppt', '.pptx']]
        else:
            raise ValueError(f"Il percorso {self.docs_dir} non è un file né una cartella valida.")

        document_futures = [process_document.remote(str(path)) for path in paths]
        documents_with_text = ray.get(document_futures)
        return documents_with_text
